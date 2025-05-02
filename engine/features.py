import struct
import time
from playsound import playsound
import os
import eel
from engine.config import *
# from engine.command import speak
from engine.TextToSpeech import TextToSpeech
from engine.SpeechToText import SpeechRecognition
import pywhatkit as kit
import sqlite3
import subprocess
import webbrowser as wb
from engine.helper import extract_yt_term, remove_words
import pvporcupine
import pyaudio
import pyautogui as autogui
from telethon.sync import TelegramClient
from telethon.tl.types import InputPeerUser
import google.generativeai as genai
from engine.config import PORCUPINE_ACCESS_KEY, GENAI_API_KEY, HUGGING_FACE_TOKEN
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from huggingface_hub import InferenceClient
import random
from timeit import default_timer as timer
import cv2

@eel.expose
def playAssistantSound():
    # Absolute path to the audio file
    music_dir = "www/assets/audio/www_assets_audio_start_sound.mp3"
    try:
        playsound(music_dir)
    except Exception as e:
        print(f"Error playing sound: {e}")

def openCommand(query: str):
    query.lower()
    query = query.replace(ASSISTANT_NAME, "")
    query = query.replace("open", "")
    
    app_name = query.strip()
    if app_name != "":
        try:
            con = sqlite3.connect("jarvis.db")
            cursor = con.cursor()
            sql_query = "SELECT path FROM sys_command WHERE name IN (?)"
            cursor.execute(sql_query, (app_name,))
            results = cursor.fetchall()
            if len(results) != 0:
                TextToSpeech(f"Opening {query}...")
                subprocess.call(["open", results[0][0]])
            elif len(results) == 0:
                sql_query = "SELECT path FROM web_command WHERE name IN (?)"
                cursor.execute(sql_query, (app_name,))
                results = cursor.fetchall()
                if len(results) != 0:
                    TextToSpeech(f"Opening {query}...")
                    wb.open(results[0][0])
                else:
                    TextToSpeech(f"Opening {query}...")
                    try:
                        os.system(f"open -a {query}")
                    except:
                        TextToSpeech("Application not found...")
            con.close()
        except:
            TextToSpeech("Something went Wrong...")

def PlayYoutube(query: str):
    search_term = extract_yt_term(query)
    if search_term != None:
        TextToSpeech(f"Playing {search_term} on YouTube...")
        kit.playonyt(search_term)
    else:
        TextToSpeech(f"Fail to recognize search term...")

def trigger_shortcut():
    """Trigger the Cmd+D shortcut."""
    autogui.keyDown("ctrl")
    autogui.press("d")
    time.sleep(2)
    autogui.keyUp("ctrl")


def hotword():
    porcupine = None
    paud = None
    audio_stream = None
    try:
        # Initialize Porcupine for hotword detection
        porcupine = pvporcupine.create(keywords=["jarvis", "alexa"], access_key=PORCUPINE_ACCESS_KEY)
        paud = pyaudio.PyAudio()

        # Open an audio stream
        audio_stream = paud.open(
            rate=porcupine.sample_rate,
            channels=1,
            format=pyaudio.paInt16,
            input=True,
            frames_per_buffer=porcupine.frame_length*2,
        )

        print("Listening for hotwords...")

        # Loop for audio stream processing
        while True:
            try:
                keyword_frame = audio_stream.read(porcupine.frame_length, exception_on_overflow=False)
                keyword_frame = struct.unpack_from("h" * porcupine.frame_length, keyword_frame)
                keyword_index = porcupine.process(keyword_frame)

                if keyword_index >= 0:
                    print("Hotword detected!")
                    trigger_shortcut()
            except IOError as e:
                print(f"Audio input error: {e}")
                continue

    except Exception as e:
        print(f"Error: {e}")
    finally:
        # Cleanup resources
        if porcupine is not None:
            porcupine.delete()
        if audio_stream is not None:
            audio_stream.close()
        if paud is not None:
            paud.terminate()

# Finding Contact in Database
def findContact(query):
    con = sqlite3.connect("jarvis.db")
    cursor = con.cursor()
    words_to_remove = [ASSISTANT_NAME, 'make', 'a', 'to', 'phone', 'call', 'send', 'message', 'wahtsapp', 'video']
    query = remove_words(query, words_to_remove)

    try:
        query = query.strip().lower()
        cursor.execute("SELECT api_id, api_hash FROM telegram WHERE LOWER(name) LIKE ? OR LOWER(name) LIKE ?", ('%' + query + '%', query + '%'))
        results = cursor.fetchall()
        api_id = results[0][0]
        api_hash = results[0][1]

        return (api_id, api_hash)
    except:
        TextToSpeech('not exist in contacts')
        return 0, 0

def telegram(user_id, user_hash, message):
    client = TelegramClient('session', API_ID, API_HASH)
    try:
        client.connect()
        if not client.is_user_authorized():
            client.send_code_request(PHONE)
            client.sign_in(PHONE, input('Enter the code: '))
        receiver = InputPeerUser(user_id, user_hash)
        client.send_message(receiver, message, parse_mode='html')
        print("Message sent successfully!")
    except Exception as e:
        print(f"An error occurred: {e}")
    finally:
        client.disconnect()

def chatBot(query: str):
    user_input = query.lower()
    genai.configure(api_key=GENAI_API_KEY)
    model = genai.GenerativeModel("gemini-1.5-flash")
    response = model.generate_content(f"'{user_input}'")
    response = response.text.replace("*", "")
    # print(response)
    TextToSpeech(response)

def configure_clients():
    """Configure Generative AI and Stable Diffusion clients."""
    genai.configure(api_key=GENAI_API_KEY)
    # client = InferenceClient(model="stabilityai/stable-diffusion-3.5-large-turbo", 
    #                          token=HUGGING_FACE_TOKEN)
    client = InferenceClient(provider="hf-inference", api_key=HUGGING_FACE_TOKEN)
    model = genai.GenerativeModel("gemini-1.5-flash-latest", 
                                  generation_config={"response_mime_type": "application/json"})
    return client, model


def generate_presentation_content(model, num_slides, topic):
    """Generate presentation content using the Generative AI model."""
    prompt = (
        f"Generate a {num_slides}-slide presentation on the topic '{topic}'. "
        "Each slide should have a 'slide_no', 'header', 'content', and 'image_prompt' section. "
        "Provide 50-60 words of content per slide. "
        "Return the output in JSON format."
    )
    raw_response = model.generate_content(prompt)
    return json.loads(raw_response.text)


def generate_image(client, prompt, img_path, index, max_retries=5):
    """Generate an image using the Stable Diffusion client with retry logic."""
    retries = 0
    while retries < max_retries:
        try:
            # output = client.text_to_image(prompt+", Ultra 4K resolution.")
            output = client.text_to_image(
                prompt+", Ultra 4K resolution.",
                model="stabilityai/stable-diffusion-xl-base-1.0",
            )
            image_path = os.path.join(img_path, f"{index}.png")
            output.save(image_path)
            return image_path
        except Exception as e:
            print(f"Error generating image: {e}. Retrying... ({retries + 1}/{max_retries})")
            retries += 1
            time.sleep(5)  # Wait before retrying
    
    raise RuntimeError(f"Failed to generate image after {max_retries} attempts.")


def create_slide(prs, slide_data, client, img_path, bg_path, slide_layout, index):
    """Create a single slide with content and an image."""
    new_slide = prs.slides.add_slide(slide_layout)
    
    # Set background
    left = top = Inches(0)
    pic = new_slide.shapes.add_picture(bg_path, left, top, width=prs.slide_width, height=prs.slide_height)
    new_slide.shapes._spTree.remove(pic._element)
    new_slide.shapes._spTree.insert(2, pic._element)
    # Add slide title
    title = new_slide.shapes.title
    title.text = slide_data['header']
    
    # Add slide content
    content_box = new_slide.placeholders[2 if int(slide_data['slide_no']) % 2 == 0 else 1]
    tf = content_box.text_frame
    tf.text = slide_data['content'].replace(".", ".\n").strip()
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(20)
    
    # Generate and add image
    image_path = generate_image(client, slide_data['image_prompt'], img_path, index)
    content_box = new_slide.placeholders[1 if int(slide_data['slide_no']) % 2 == 0 else 2]
    new_slide.shapes.add_picture(image_path, content_box.left, content_box.top, content_box.width, content_box.height)
    TextToSpeech(f"Slide {slide_data['slide_no']} completed successfully.")
    print(f"Slide {slide_data['slide_no']} completed successfully.")


def generate_presentation(response, client, img_path, bg_folder):
    """Generate the PowerPoint presentation from the response."""
    prs = Presentation()
    r1 = random.randint(1, 5)
    bg_path = os.path.join(bg_folder, f"{r1}.jpg")
    slide_layout = prs.slide_layouts[3]
    
    for i, slide_data in enumerate(response['presentation']):
        create_slide(prs, slide_data, client, img_path, bg_path, slide_layout, i)
        time.sleep(10)
    
    return prs


def presentation(topic, num=2):
    client, model = configure_clients()
    
    # User inputs
    img_path = "assets/images/"
    bg_folder = "assets/background/"
    output_path = "assets/ppts/test.pptx"
    
    start = timer()
    
    try:
        response = generate_presentation_content(model, num, topic)
        prs = generate_presentation(response, client, img_path, bg_folder)
        prs.save(output_path)
        print(f"Presentation saved to {output_path}")
        TextToSpeech("Opening ppt...")
        os.system(f'open "{output_path}"')
    except Exception as e:
        prs.save(output_path)
        print(f"An error occurred: {e}")
    finally:
        end = timer()
        print(f"Time taken: {(end - start) / 60:.2f} minutes")

def generateImage(topic: str):
    TextToSpeech(f"Generating an image of {topic}...")
    client = InferenceClient("stabilityai/stable-diffusion-3.5-large-turbo", token=HUGGING_FACE_TOKEN)
    # output is a PIL.Image object
    output = client.text_to_image(topic)
    topic = topic.replace(" ", "_")
    output.save(f"assets/image/{topic}.png")  # Save the image
    output.show()


def takePic():
    camera_port = 0
    camera = cv2.VideoCapture(camera_port)
    time.sleep(1.5)  # If you don't wait, the image will be dark
    return_value, image = camera.read()
    print(return_value)
    cv2.imwrite("Camera/opencv.png", image)
    del(camera)  # so that others can use the camera as soon as possible
    os.system(f'open "Camera/opencv.png"')


if __name__ == "__main__":
    print("Hello")